import streamlit as st
import pandas as pd
import io
import re
from PIL import Image
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.drawing.image import Image as XLImage

# 新增讀檔套件 (如果尚未安裝，請依之前的教學安裝)
import pdfplumber
from pptx import Presentation

# --- 1. 頁面配置 ---
st.set_page_config(page_title="營造履歷智慧填表系統 v8.0", layout="wide", page_icon="🏗️")

# --- 2. CSS 樣式 (維持風格) ---
st.markdown("""
    <style>
    :root { --main-yellow: #FFB81C; --accent-orange: #FF4438; --dark-grey: #2D2926; }
    .stApp { background-color: #f4f6f9; }
    h1, h2, h3, label { color: var(--dark-grey) !important; font-weight: bold !important; font-family: '微軟正黑體', sans-serif; }
    .stButton>button { 
        background-color: var(--main-yellow); color: var(--dark-grey); 
        border: none; width: 100%; border-radius: 8px; font-size: 18px; font-weight: bold; padding: 12px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }
    div[data-testid="stExpander"] { background-color: white; border-radius: 10px; box-shadow: 0 2px 5px rgba(0,0,0,0.05); }
    </style>
    """, unsafe_allow_html=True)

# --- 3. 智慧提取函式 ---
def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_construction_data(text):
    data = {}
    name_match = re.search(r"(\S*新建工程|\S*大樓工程)", text)
    if name_match: data["project_name"] = name_match.group(1)

    area_match = re.search(r"基地面積\D*([\d,]+\.?\d*)", text)
    if area_match:
        try: data["site_area"] = float(area_match.group(1).replace(",", ""))
        except: pass

    fa_match = re.search(r"(總樓地板|總樓地|總建坪)\D*([\d,]+\.?\d*)", text)
    if fa_match:
        try: data["total_floor_area"] = float(fa_match.group(2).replace(",", ""))
        except: pass

    up_match = re.search(r"地上\D*(\d+)", text)
    down_match = re.search(r"地下\D*(\d+)", text)
    if up_match: data["floors_up"] = int(up_match.group(1))
    if down_match: data["floors_down"] = int(down_match.group(1))

    depth_match = re.search(r"(開挖深度|GL-)\D*([\d,]+\.?\d*)", text)
    if depth_match:
        try: data["excavation_depth"] = float(depth_match.group(2).replace(",", ""))
        except: pass

    if "逆打" in text: data["const_method"] = "逆打工法 (Top-Down)"
    elif "雙順打" in text: data["const_method"] = "雙順打工法"
    elif "順打" in text: data["const_method"] = "順打工法 (Bottom-Up)"

    if "SRC" in text: data["struct_above"] = "SRC (鋼骨鋼筋混凝土)"
    elif "SC" in text: data["struct_above"] = "SC (鋼骨)"
    elif "RC" in text: data["struct_above"] = "RC (鋼筋混凝土)"

    return data

# --- 4. 初始化 Session State (設定為空白預設值) ---
# 這裡全部設為空字串或 0，讓介面乾淨
default_values = {
    "project_name": "",
    "project_loc": "",
    "client_name": "",
    "architect_name": "",
    "contract_date": "",
    "contract_cost": "",
    "floors_up": 0,
    "floors_down": 0,
    "site_area": 0.0,
    "total_floor_area": 0.0,
    "building_height": 0.0,
    "excavation_depth": 0.0,
    "const_method": "請選擇...",
    "struct_above": "請選擇...",
    "struct_below": "請選擇...",
    "foundation_type": "請選擇...",
    "b_type": "請選擇...",
    "retain_sys": "請選擇...",
    "wall_sys": "請選擇..."
}

for key, val in default_values.items():
    if key not in st.session_state:
        st.session_state[key] = val

# --- 5. 介面設計 ---

st.title("🏗️ 營造履歷智慧填表系統 v8.0")

# === 檔案上傳區 ===
with st.expander("📂 智慧匯入 (拖曳 PDF/PPT 檔案到此)", expanded=True):
    col_up1, col_up2 = st.columns([2, 1])
    with col_up1:
        uploaded_doc = st.file_uploader("若有標案簡報，可直接拖曳至此自動填寫", type=["pdf", "pptx"])
    with col_up2:
        st.write("") # Spacer
        st.write("")
        if uploaded_doc is not None:
            if st.button("🚀 開始分析檔案", type="primary"):
                with st.spinner("正在讀取檔案..."):
                    try:
                        raw_text = ""
                        if uploaded_doc.name.endswith(".pdf"):
                            raw_text = extract_text_from_pdf(uploaded_doc)
                        elif uploaded_doc.name.endswith(".pptx"):
                            raw_text = extract_text_from_ppt(uploaded_doc)
                        
                        extracted_data = parse_construction_data(raw_text)
                        
                        if extracted_data:
                            for k, v in extracted_data.items():
                                st.session_state[k] = v
                            st.success(f"✅ 自動填入 {len(extracted_data)} 欄位！")
                            st.rerun() # 重新整理頁面以顯示資料
                        else:
                            st.warning("⚠️ 未偵測到關鍵字，請手動輸入")
                    except Exception as e:
                        st.error(f"解析失敗：{e}")

st.markdown("---")

# === 填表區 (加入 placeholder 提示) ===
tab1, tab2, tab3 = st.tabs(["📝 基本資料與規格", "🖼️ 圖片與敘述", "📊 導出 Excel"])

with tab1:
    st.subheader("1. 專案基本資料")
    c1, c2, c3 = st.columns(3)
    with c1:
        st.text_input("專案名稱", key="project_name", placeholder="例：信義區A1新建工程") 
        st.text_input("工程地點", key="project_loc", placeholder="例：台北市信義區")
    with c2:
        st.text_input("業主名稱", key="client_name", placeholder="例：XX建設股份有限公司")
        st.text_input("設計單位/建築師", key="architect_name", placeholder="例：OOO建築師事務所")
    with c3:
        st.text_input("完工年份", key="contract_date", placeholder="例：2023.05 - 2025.12")
        st.text_input("工程造價 (億元)", key="contract_cost", placeholder="例：15.5")

    st.subheader("2. 建築規模")
    
    # 輔助函式：處理下拉選單，讓預設值正確顯示
    def get_index(options, key):
        current_val = st.session_state[key]
        if current_val in options:
            return options.index(current_val)
        return 0

    col_b1, col_b2, col_b3, col_b4 = st.columns(4)
    with col_b1:
        opts_type = ["請選擇...", "住宅大樓", "商辦大樓", "飯店", "廠房", "公共工程"]
        st.selectbox("建物類型", opts_type, index=get_index(opts_type, "b_type"), key="b_type")
    with col_b2:
        opts_struct = ["請選擇...", "SC (鋼骨)", "SRC (鋼骨鋼筋混凝土)", "RC (鋼筋混凝土)", "SS (純鋼構)"]
        st.selectbox("地上結構", opts_struct, index=get_index(opts_struct, "struct_above"), key="struct_above")
    with col_b3:
        opts_struct_down = ["請選擇...", "RC (鋼筋混凝土)", "SRC (鋼骨鋼筋混凝土)"]
        st.selectbox("地下結構", opts_struct_down, index=get_index(opts_struct_down, "struct_below"), key="struct_below")
    with col_b4:
        opts_found = ["請選擇...", "筏式基礎", "筏式基礎+基樁", "獨立基腳"]
        st.selectbox("基礎型式", opts_found, index=get_index(opts_found, "foundation_type"), key="foundation_type")

    col_d1, col_d2, col_d3 = st.columns(3)
    with col_d1:
        st.number_input("地上層數 (F)", min_value=0, key="floors_up", help="輸入 0 表示未定")
        st.number_input("地下層數 (B)", min_value=0, key="floors_down")
    with col_d2:
        st.number_input("基地面積 (m²)", key="site_area")
        st.number_input("總樓地板面積 (m²)", key="total_floor_area")
    with col_d3:
        st.number_input("建築高度 (m)", key="building_height")
        st.number_input("開挖深度 (m)", key="excavation_depth")

    st.subheader("3. 關鍵工法")
    c_m1, c_m2, c_m3 = st.columns(3)
    with c_m1:
        opts_method = ["請選擇...", "逆打工法 (Top-Down)", "順打工法 (Bottom-Up)", "雙順打工法"]
        st.selectbox("主體施工工法", opts_method, index=get_index(opts_method, "const_method"), key="const_method")
    with c_m2:
        opts_retain = ["請選擇...", "連續壁+鋼支柱(逆打)", "連續壁+內支撐", "地錨工法", "鋼板樁"]
        st.selectbox("擋土支撐系統", opts_retain, index=get_index(opts_retain, "retain_sys"), key="retain_sys")
    with c_m3:
        opts_wall = ["請選擇...", "玻璃帷幕", "石材吊掛", "鋁板", "二丁掛"]
        st.selectbox("外牆工法", opts_wall, index=get_index(opts_wall, "wall_sys"), key="wall_sys")

with tab2:
    st.header("工程特色與圖片")
    col_text1, col_text2 = st.columns(2)
    with col_text1:
        features = st.text_area("✨ 工程特色 (條列式)", placeholder="1. 採用特殊工法...\n2. 獲得綠建築標章...", height=200)
    with col_text2:
        challenges = st.text_area("🧗 施工挑戰 (條列式)", placeholder="1. 鄰近捷運監測...\n2. 基地狹小...", height=200)

    uploaded_img = st.file_uploader("上傳完工照 (JPG/PNG)", type=['jpg', 'png', 'jpeg'])
    if uploaded_img:
        st.image(uploaded_img, width=400, caption="預覽圖片")

with tab3:
    st.header("導出 Excel 履歷")
    
    def generate_excel():
        wb = Workbook()
        ws = wb.active
        ws.title = "專案履歷表"
        
        # 簡易檢查：如果沒填資料，提醒使用者
        p_name = st.session_state.project_name if st.session_state.project_name else "未命名專案"
        
        # 樣式與欄寬設定
        border_style = Side(border_style="thin", color="000000")
        full_border = Border(left=border_style, right=border_style, top=border_style, bottom=border_style)
        fill_header = PatternFill(start_color="2D2926", end_color="2D2926", fill_type="solid")
        fill_sub_header = PatternFill(start_color="FFB81C", end_color="FFB81C", fill_type="solid")
        fill_light = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
        font_title = Font(name='微軟正黑體', size=16, bold=True, color="FFFFFF")
        font_sub = Font(name='微軟正黑體', size=12, bold=True)
        font_label = Font(name='微軟正黑體', size=11, bold=True)
        font_val = Font(name='微軟正黑體', size=11)

        ws.column_dimensions['A'].width = 15
        ws.column_dimensions['B'].width = 25
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 25

        # 標題
        ws.merge_cells('A1:D1')
        ws['A1'] = p_name
        ws['A1'].fill = fill_header
        ws['A1'].font = font_title
        ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[1].height = 40

        def write_row(r, l1, v1, l2, v2):
            ws[f'A{r}'] = l1
            ws[f'B{r}'] = v1
            ws[f'C{r}'] = l2
            ws[f'D{r}'] = v2
            for c in ['A','C']: 
                ws[f'{c}{r}'].fill = fill_light
                ws[f'{c}{r}'].font = font_label
            for c in ['B','D']: ws[f'{c}{r}'].font = font_val
            for c in ['A','B','C','D']: 
                ws[f'{c}{r}'].border = full_border
                ws[f'{c}{r}'].alignment = Alignment(vertical='center', wrap_text=True)

        # 寫入資料 (從 session_state 讀取)
        ss = st.session_state
        write_row(2, "工程地點", ss.project_loc, "完工年份", ss.contract_date)
        write_row(3, "業主單位", ss.client_name, "設計單位", ss.architect_name)
        cost_str = f"{ss.contract_cost} 億元" if ss.contract_cost else ""
        write_row(4, "工程造價", cost_str, "建物用途", ss.b_type)

        # 分隔
        ws.merge_cells('A5:D5')
        ws['A5'] = "建築規模與技術規格"
        ws['A5'].fill = fill_sub_header
        ws['A5'].font = font_sub
        ws['A5'].alignment = Alignment(horizontal='center')
        ws['A5'].border = full_border

        struct_str = f"地上:{ss.struct_above} / 地下:{ss.struct_below}"
        floor_str = f"{ss.floors_up}F / {ss.floors_down}B (高 {ss.building_height}m)"
        area_str = f"基地:{ss.site_area:,.0f} / 總樓:{ss.total_floor_area:,.0f} m²"
        excav_str = f"{ss.const_method} / GL-{ss.excavation_depth}m"

        write_row(6, "樓層/高度", floor_str, "結構系統", struct_str)
        write_row(7, "面積資訊", area_str, "基礎型式", ss.foundation_type)
        write_row(8, "施工工法", excav_str, "擋土系統", ss.retain_sys)
        write_row(9, "外牆系統", ss.wall_sys, "其他", "")

        # 特色與圖片 (略為簡化，與上版相同邏輯)
        ws.merge_cells('A10:D10'); ws['A10'] = "工程特色"; ws['A10'].fill = fill_sub_header; ws['A10'].font = font_sub; ws['A10'].border = full_border
        ws.merge_cells('A11:D11'); ws['A11'] = features if features else "(無)"; ws['A11'].alignment = Alignment(wrap_text=True, vertical='top'); ws['A11'].border = full_border; ws.row_dimensions[11].height = 60
        
        ws.merge_cells('A12:D12'); ws['A12'] = "施工挑戰"; ws['A12'].fill = fill_sub_header; ws['A12'].font = font_sub; ws['A12'].border = full_border
        ws.merge_cells('A13:D13'); ws['A13'] = challenges if challenges else "(無)"; ws['A13'].alignment = Alignment(wrap_text=True, vertical='top'); ws['A13'].border = full_border; ws.row_dimensions[13].height = 60
        
        ws.merge_cells('A14:D14'); ws['A14'] = "專案照片"; ws['A14'].fill = fill_sub_header; ws['A14'].font = font_sub; ws['A14'].alignment = Alignment(horizontal='center'); ws['A14'].border = full_border

        if uploaded_img:
            img_io = io.BytesIO(uploaded_img.getvalue())
            img = XLImage(img_io)
            img.width = 400; img.height = 300
            ws.add_image(img, 'A15')
            ws.row_dimensions[15].height = 230
        else:
            ws.merge_cells('A15:D15')
            ws['A15'] = "(無照片)"
            ws['A15'].alignment = Alignment(horizontal='center', vertical='center')
            ws.row_dimensions[15].height = 50

        out_buffer = io.BytesIO()
        wb.save(out_buffer)
        return out_buffer.getvalue()

    if st.button("生成並下載 Excel", type="primary"):
        xlsx_data = generate_excel()
        p_name = st.session_state.project_name if st.session_state.project_name else "Project"
        st.download_button(
            label="📥 點擊下載",
            data=xlsx_data,
            file_name=f"{p_name}_履歷表.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )